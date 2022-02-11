# https://python-pptx.readthedocs.io/en/latest/index.html#
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

prs = Presentation("samples/template.pptx")
# 첫번째 슬라이드 제목/부제목 넣기
slide1 = prs.slides[0]
slide1.placeholders[0].text = "Hello world!"
slide1.placeholders[0].text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
slide1.placeholders[1].text = "powered by Python"
slide1.placeholders[1].text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# 두번째 슬라이드에 제목, 테이블 채우기
slide2 = prs.slides[1]
slide2.placeholders[0].text = "Table example"
table = slide2.shapes[1].table
table.cell(0, 0).text = "Column 1"
table.cell(0, 1).text = "Column 2"
table.cell(0, 2).text = "Column 3"
table.cell(0, 3).text = "Column 4"
table.cell(0, 4).text = "Column 5"

table.cell(1, 0).text = "Value 1"
table.cell(2, 3).text = "Value 2,3"

# 세번째 슬라이드 새로 추가
slide3 = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
pic = slide3.shapes.add_picture("samples/python-logo.png", Inches(1), Inches(1))

prs.save("samples/result.pptx")